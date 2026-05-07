"""
Brain-Viz presentation — 20 slides
Dark neuroscience theme: deep navy + electric blue + teal + gold
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ─────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x11, 0x17)   # very dark navy
BG_MID    = RGBColor(0x0F, 0x16, 0x29)   # dark blue
BG_CARD   = RGBColor(0x15, 0x1E, 0x35)   # card surface
ACCENT    = RGBColor(0x3A, 0x86, 0xFF)   # electric blue
TEAL      = RGBColor(0x2E, 0xC4, 0xB6)   # teal
GOLD      = RGBColor(0xFF, 0xD1, 0x66)   # gold
PINK      = RGBColor(0xF7, 0x25, 0x85)   # magenta
GREEN     = RGBColor(0x06, 0xD6, 0xA0)   # mint green
ORANGE    = RGBColor(0xFF, 0x6B, 0x35)   # orange
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_HI   = RGBColor(0xE2, 0xE8, 0xF8)   # bright text
TEXT_MID  = RGBColor(0x7E, 0x95, 0xCC)   # muted blue-white
TEXT_DIM  = RGBColor(0x44, 0x50, 0x70)   # dim text
BLANK     = RGBColor(0xFF, 0xFF, 0xFF)

def blank_slide():
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color=BG_DARK):
    from pptx.util import Emu
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, alpha=None, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=18, bold=False, italic=False,
        color=TEXT_HI, align=PP_ALIGN.LEFT,
        wrap=True, valign=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return tb

def title_bar(slide, title, accent_color=ACCENT):
    """Coloured left-side accent bar + title text"""
    rect(slide, 0.45, 0.3, 0.07, 0.75, accent_color)
    txt(slide, title, 0.65, 0.22, 11.5, 0.9,
        size=32, bold=True, color=TEXT_HI, align=PP_ALIGN.LEFT)

def card(slide, x, y, w, h, color=BG_CARD):
    shape = rect(slide, x, y, w, h, color)
    return shape

def bullet_block(slide, items, x, y, w, h, size=15, color=TEXT_HI,
                 accent=ACCENT, spacing=0.52):
    """Draw bullet items with a colored dot"""
    cy = y
    for item in items:
        rect(slide, x, cy + 0.10, 0.06, 0.06, accent)
        txt(slide, item, x + 0.18, cy, w - 0.18, spacing,
            size=size, color=color, align=PP_ALIGN.LEFT)
        cy += spacing

def stat_box(slide, number, label, x, y, w=2.4, accent=ACCENT):
    card(slide, x, y, w, 1.35)
    rect(slide, x, y, w, 0.07, accent)
    txt(slide, number, x, y+0.12, w, 0.7, size=34, bold=True,
        color=accent, align=PP_ALIGN.CENTER)
    txt(slide, label,  x, y+0.80, w, 0.45, size=13,
        color=TEXT_MID, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1  —  TITLE
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)

# Left gradient panel
rect(s, 0, 0, 5.5, 7.5, BG_MID)
rect(s, 0, 0, 0.12, 7.5, ACCENT)

# Decorative neural-network dots (circles)
for cx, cy, cr, col in [
    (1.2,1.0,0.18,ACCENT),(2.8,0.6,0.12,TEAL),(0.7,2.5,0.10,PINK),
    (3.5,1.8,0.15,GOLD), (1.8,3.2,0.10,GREEN),(4.2,0.9,0.08,ACCENT),
    (2.1,1.5,0.08,TEAL), (4.8,2.4,0.12,PINK),
]:
    sh = s.shapes.add_shape(9, Inches(cx-cr), Inches(cy-cr), Inches(cr*2), Inches(cr*2))
    sh.fill.solid(); sh.fill.fore_color.rgb = col
    sh.line.fill.background()

txt(s, "Brain-Viz", 0.5, 1.8, 4.6, 1.3, size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txt(s, "Interactive 3D Visual Analytics", 0.5, 3.0, 4.6, 0.6, size=22, italic=True, color=TEAL, align=PP_ALIGN.LEFT)
txt(s, "Human Connectomics Visualization", 0.5, 3.55, 4.6, 0.5, size=17, color=TEXT_MID, align=PP_ALIGN.LEFT)

rect(s, 0.5, 4.2, 3.8, 0.04, ACCENT)

txt(s, "Nishant  ·  2025EET2734", 0.5, 4.4, 4.6, 0.4, size=14, color=TEXT_MID, align=PP_ALIGN.LEFT)
txt(s, "Aryan  ·  2025EET2486",   0.5, 4.8, 4.6, 0.4, size=14, color=TEXT_MID, align=PP_ALIGN.LEFT)
txt(s, "EET — Data Visualization  |  2025", 0.5, 5.6, 4.6, 0.35, size=12, color=TEXT_DIM, align=PP_ALIGN.LEFT)

# Right panel — abstract brain grid
for row in range(8):
    for col in range(9):
        import random; random.seed(row*9+col)
        if random.random() > 0.55:
            alpha_node = [ACCENT,TEAL,PINK,GOLD,GREEN][random.randint(0,4)]
            cx2 = 5.9 + col*0.82
            cy2 = 0.5 + row*0.82
            r2  = random.choice([0.08,0.11,0.14,0.09])
            sh2 = s.shapes.add_shape(9, Inches(cx2-r2), Inches(cy2-r2), Inches(r2*2), Inches(r2*2))
            sh2.fill.solid(); sh2.fill.fore_color.rgb = alpha_node
            sh2.line.fill.background()

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2  —  INTRODUCTION & MOTIVATION
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Introduction & Motivation", ACCENT)

# Left column
card(s, 0.45, 1.25, 5.9, 5.6, BG_CARD)
txt(s, "Why Brain Connectivity?", 0.65, 1.35, 5.5, 0.5,
    size=19, bold=True, color=ACCENT)
items = [
    "The human brain contains ~86 billion neurons forming trillions of synaptic connections",
    "Structural connectivity (connectome) governs cognition, emotion, and behaviour",
    "Neurological disorders — Alzheimer's, autism, schizophrenia — manifest as connectivity disruptions",
    "Visual analytics can reveal patterns invisible in raw matrices",
    "Clinicians and researchers need intuitive tools, not just data dumps",
]
bullet_block(s, items, 0.65, 1.95, 5.5, 5.0, size=14, spacing=0.98)

# Right column — 3 stat cards
for val, lbl, acc, cx in [
    ("86B",  "Neurons in\nhuman brain",        ACCENT, 6.65),
    ("100T", "Synaptic\nconnections",           TEAL,   6.65),
    ("50M+", "People with\nneurological disease", GOLD, 6.65),
]:
    pass  # built below with correct y offsets

for i, (val, lbl, acc) in enumerate([
    ("86 B",  "Neurons in\nhuman brain", ACCENT),
    ("100 T", "Synaptic\nconnections", TEAL),
    ("50 M+", "Global neurological\ndisease burden", GOLD),
]):
    stat_box(s, val, lbl, 6.65, 1.3 + i*1.85, 2.5, acc)

card(s, 6.5, 6.2, 6.5, 1.0, RGBColor(0x1C,0x28,0x44))
txt(s, "💡  Our Goal: A zero-install, web-based visual analytics tool that makes brain "
       "connectivity exploration accessible to researchers AND laypersons alike.",
    6.65, 6.25, 6.3, 0.85, size=12, color=TEXT_HI, italic=True)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3  —  PROBLEM DEFINITION
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Problem Definition", PINK)

cols = [
    ("🧩  The Challenge",  PINK,  0.45, [
        "Brain networks involve 100s of regions with 1000s of connections",
        "Existing tools (BrainNet Viewer, Connectome Workbench) require installation of heavy software",
        "No unified tool offers both 3D spatial AND 2D topological views simultaneously",
        "Pathfinding between brain regions is not natively supported",
        "Comparing healthy vs. disease-affected brains requires manual analysis",
    ]),
    ("✅  Our Solution", GREEN, 7.0, [
        "Lightweight web app — runs entirely in the browser, zero install",
        "Dual-view: 3D WebGL spatial brain + D3.js 2D chord topology",
        "Brushing & linking between views for cross-view exploration",
        "Dijkstra pathfinding to trace strongest structural pathways",
        "Side-by-side comparison of two connectomes with structural diff",
    ]),
]
for title, acc, cx, pts in cols:
    card(s, cx, 1.2, 5.9, 5.9, BG_CARD)
    rect(s, cx, 1.2, 5.9, 0.07, acc)
    txt(s, title, cx+0.2, 1.28, 5.5, 0.55, size=18, bold=True, color=acc)
    bullet_block(s, pts, cx+0.2, 1.95, 5.5, 5.0, size=13.5, spacing=1.0)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4  —  RELATED WORK
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Related Work", TEAL)

tools = [
    ("BrainNet\nViewer",      "Xia et al.", ACCENT,
     "MATLAB-based 3D brain network visualizer.\nWidely used in neuroscience research.\nLimitation: Requires MATLAB licence, not web-accessible."),
    ("Connectome\nWorkbench", "HCP",        TEAL,
     "Full-featured medical imaging platform.\nSupports multi-modal neuroimaging data.\nLimitation: Complex interface, not suited for exploratory analytics."),
    ("Brain\nVoyager",        "Goebel",     GOLD,
     "Commercial fMRI analysis software.\nRich feature set for clinical research.\nLimitation: Paid licence, no web deployment."),
    ("D3 / WebGL\nLibraries", "Open Source",PINK,
     "General-purpose graph & 3D rendering.\nUsed in several neuroscience dashboards.\nLimitation: No domain-specific brain atlas integration."),
]
for i,(name, src, acc, desc) in enumerate(tools):
    cx = 0.45 + i*3.22
    card(s, cx, 1.2, 3.0, 4.2, BG_CARD)
    rect(s, cx, 1.2, 3.0, 0.07, acc)
    txt(s, name, cx+0.15, 1.3, 2.7, 0.8, size=16, bold=True, color=acc)
    txt(s, src,  cx+0.15, 2.0, 2.7, 0.35, size=11, italic=True, color=TEXT_MID)
    txt(s, desc, cx+0.15, 2.4, 2.7, 2.8, size=12, color=TEXT_HI)

card(s, 0.45, 5.6, 12.4, 1.5, RGBColor(0x1C,0x28,0x44))
rect(s, 0.45, 5.6, 12.4, 0.07, ACCENT)
txt(s, "How Brain-Viz differs:  Web-native  ·  Zero install  ·  "
       "Dual-view (3D+2D)  ·  Pathfinding  ·  Real-data pipeline  ·  Disease comparison",
    0.65, 5.72, 12.0, 0.7, size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "Brain-Viz is the only tool combining all these features in a single browser-based platform.",
    0.65, 6.3, 12.0, 0.5, size=13, color=TEXT_MID, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5  —  DATA SOURCES
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Data Sources & Acquisition", GOLD)

sources = [
    ("HCP",   "Human Connectome\nProject",  ACCENT, "Healthy adult brains\nDTI structural connectivity\n100–500 parcellation regions\nSpatial MNI coordinates"),
    ("ABIDE", "Autism Brain\nImaging Exchange",TEAL, "Autism vs. neurotypical\nfMRI resting-state data\nCPAC preprocessed\nOpen access, 539 subjects"),
    ("ADNI",  "Alzheimer's Disease\nNeuroimaging",GOLD,"Alzheimer's progression\nDTI + fMRI longitudinal\nFree academic access\nRegistration required"),
    ("nilearn","Python Atlas\nLibrary",      PINK, "Auto-downloads atlases\nSchaefer 100/200/400\nDestrieux cortical parcels\nNo manual files needed"),
]
for i,(abbr,name,acc,desc) in enumerate(sources):
    cx = 0.45 + i*3.22
    card(s, cx, 1.2, 3.0, 3.8, BG_CARD)
    rect(s, cx, 1.2, 3.0, 0.07, acc)
    txt(s, abbr, cx+0.12, 1.3,  2.76, 0.55, size=22, bold=True, color=acc)
    txt(s, name, cx+0.12, 1.82, 2.76, 0.65, size=12, italic=True, color=TEXT_MID)
    txt(s, desc, cx+0.12, 2.5,  2.76, 2.3,  size=12, color=TEXT_HI)

# Pipeline
card(s, 0.45, 5.2, 12.4, 2.0, BG_MID)
txt(s, "Data Pipeline", 0.65, 5.28, 3.0, 0.4, size=14, bold=True, color=GOLD)
steps = ["Download\nfMRI/DTI","Atlas\nregistration","Connectivity\nmatrix","Threshold\n& normalise","process_hcp.py\n→ JSON","Brain-Viz\nvisualization"]
for i,st in enumerate(steps):
    cx3 = 0.65 + i*2.08
    card(s, cx3, 5.72, 1.85, 1.3, BG_CARD)
    txt(s, str(i+1), cx3+0.1, 5.76, 0.35, 0.35, size=14, bold=True, color=GOLD)
    txt(s, st, cx3+0.1, 6.12, 1.65, 0.8, size=11, color=TEXT_HI)
    if i<5:
        txt(s, "→", cx3+1.9, 6.15, 0.25, 0.4, size=18, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6  —  SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "System Architecture", ACCENT)

layers = [
    ("Data Layer",       BG_CARD,  GOLD,   0.45, 1.2, 12.4, 1.1,
     "process_hcp.py  ·  nilearn atlas fetch  ·  ABIDE / HCP / ADNI CSVs  →  connectome.json"),
    ("State Layer",      BG_CARD,  TEAL,   0.45, 2.5, 12.4, 1.1,
     "store.js  —  Reactive event bus  ·  filter:changed  ·  path:found  ·  region:selected  ·  dimOpacity"),
    ("Algorithm Layer",  BG_CARD,  ACCENT, 0.45, 3.8, 12.4, 1.1,
     "algorithms.js  —  Dijkstra shortest path  ·  Louvain community detection  ·  Global efficiency  ·  Edge filtering"),
    ("Visualization",    BG_CARD,  PINK,   0.45, 5.1, 12.4, 1.1,
     "brain3d.js  (Three.js r128 WebGL)  ·  graph2d.js  (D3.js v7 chord + circular)  ·  Brushing & linking"),
    ("UI / UX Layer",    BG_CARD,  GREEN,  0.45, 6.4, 12.4, 0.75,
     "main.js  ·  Landing page  ·  Dataset manager  ·  Compare page  ·  Region info panel"),
]
for lbl, bgc, acc, cx, cy, w, h, desc in layers:
    card(s, cx, cy, w, h, bgc)
    rect(s, cx, cy, 0.07, h, acc)
    txt(s, lbl, cx+0.2, cy+0.1, 2.2, 0.5, size=14, bold=True, color=acc)
    txt(s, desc, cx+2.5, cy+0.22, 9.7, 0.6, size=12.5, color=TEXT_HI)
    if cy < 6.4:
        txt(s, "▼", 6.5, cy+h, 0.3, 0.35, size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7  —  ALGORITHM: LOUVAIN COMMUNITY DETECTION
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Algorithm 1: Louvain Community Detection", TEAL)

card(s, 0.45, 1.2, 6.2, 5.9, BG_CARD)
txt(s, "How It Works", 0.65, 1.3, 5.8, 0.45, size=18, bold=True, color=TEAL)
steps7 = [
    "Assign each node its own community initially",
    "For each node, compute modularity gain of moving it to each neighbour's community",
    "Move node to the community giving the highest gain (if positive)",
    "Repeat until no single move improves modularity",
    "Aggregate communities into super-nodes and repeat",
]
bullet_block(s, steps7, 0.65, 1.85, 5.8, 5.0, size=14, accent=TEAL, spacing=0.98)

txt(s, "Modularity Q  =  Σ [ Aᵢⱼ  −  kᵢkⱼ / 2m ]  ×  δ(cᵢ, cⱼ) / 2m",
    0.65, 5.8, 5.8, 0.7, size=13, color=GOLD, italic=True)

# Right: result diagram (coloured community blocks)
card(s, 6.95, 1.2, 5.9, 5.9, BG_CARD)
txt(s, "Result in Brain-Viz", 7.15, 1.3, 5.5, 0.45, size=18, bold=True, color=TEAL)
comms = [
    ("Default Mode",      ACCENT, "17 nodes  ·  medial prefrontal, PCC, angular gyrus"),
    ("Visual",            TEAL,   "28 nodes  ·  primary/extrastriate visual cortex"),
    ("Somatomotor",       GOLD,   "16 nodes  ·  pre/postcentral gyrus, SMA"),
    ("Frontoparietal",    PINK,   "21 nodes  ·  DLPFC, IPS, lateral prefrontal"),
    ("Limbic",            GREEN,  " 5 nodes  ·  hippocampus, amygdala, OFC"),
    ("Dorsal Attention",  ORANGE, "13 nodes  ·  FEF, IPS, MT+"),
]
for i,(name,col,desc) in enumerate(comms):
    cy7 = 1.9 + i*0.78
    rect(s, 7.15, cy7+0.08, 0.25, 0.25, col)
    txt(s, name, 7.55, cy7, 2.5, 0.4, size=13, bold=True, color=col)
    txt(s, desc, 7.55, cy7+0.36, 5.1, 0.34, size=11, color=TEXT_MID)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 8  —  ALGORITHM: DIJKSTRA PATHFINDING
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Algorithm 2: Dijkstra Shortest Path (Strongest Pathway)", GOLD)

card(s, 0.45, 1.2, 6.2, 5.9, BG_CARD)
txt(s, "Implementation", 0.65, 1.3, 5.8, 0.45, size=18, bold=True, color=GOLD)
pts8 = [
    "Edge cost = 1 / weight  (high connectivity = cheap to traverse)",
    "MinHeap priority queue for O((V+E) log V) performance",
    "Predecessor array enables path reconstruction",
    "Path nodes highlighted gold; animated particles travel along path",
    "Click node 1 (source) → click node 2 (destination) → path shown",
    "Displays: hop count, average edge weight, region labels",
]
bullet_block(s, pts8, 0.65, 1.85, 5.8, 5.0, size=14, accent=GOLD, spacing=0.85)

card(s, 0.65, 5.55, 5.8, 1.3, RGBColor(0x1C,0x28,0x44))
txt(s, "dist[src]=0  →  pq.push(src,0)  →  while pq: u=pq.pop()  →  "
       "for v in adj[u]: cost=dist[u]+1/w  →  if cost<dist[v]: relax",
    0.8, 5.65, 5.5, 1.1, size=11.5, color=TEAL, italic=True)

card(s, 6.95, 1.2, 5.9, 5.9, BG_CARD)
txt(s, "Interactive Demo Flow", 7.15, 1.3, 5.5, 0.45, size=18, bold=True, color=GOLD)

demo_steps = [
    (ACCENT, "1  Click source region",   "Node turns gold, 'Select destination' hint appears"),
    (GOLD,   "2  Click destination",      "Dijkstra runs on adjacency map in <5 ms"),
    (TEAL,   "3  Path renders",           "Gold line + glow drawn through intermediate nodes"),
    (PINK,   "4  Particles animate",      "14 glowing dots travel the path direction"),
    (GREEN,  "5  Info bar updates",        "Shows: L-PreCentral → R-STG · Hops: 3 · Avg: 0.62"),
    (ORANGE, "6  Click third node",        "Resets selection for new pathfinding query"),
]
for i,(acc,hd,desc) in enumerate(demo_steps):
    cy8 = 1.9 + i*0.75
    rect(s, 7.15, cy8+0.08, 0.08, 0.28, acc)
    txt(s, hd,   7.35, cy8,      4.8, 0.38, size=13, bold=True, color=acc)
    txt(s, desc, 7.35, cy8+0.38, 5.2, 0.32, size=11, color=TEXT_MID)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 9  —  ALGORITHM: CONNECTIVITY MATRIX PROCESSING
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Algorithm 3: Connectivity Matrix Processing Pipeline", PINK)

steps9 = [
    ("Input",      ACCENT, "Raw n×n connectivity matrix from DTI tractography or fMRI correlation"),
    ("Threshold",  TEAL,   "Keep only edges with weight > τ (default 0.10), discard self-loops"),
    ("Normalise",  GOLD,   "Min-max normalise all surviving weights to [0, 1]"),
    ("Cap",        PINK,   "Keep only the top-K edges (default K=700) for browser performance"),
    ("Betweenness",GREEN,  "Compute betweenness centrality via Brandes algorithm (O(VE) approximated)"),
    ("Community",  ORANGE, "Run Louvain community detection on thresholded graph"),
    ("JSON output",ACCENT, "Serialise nodes, edges, communities, metadata → connectome.json"),
]
for i,(lbl,acc,desc) in enumerate(steps9):
    cy9 = 1.25 + i*0.82
    card(s, 0.45, cy9, 12.4, 0.72, BG_CARD)
    rect(s, 0.45, cy9, 0.07, 0.72, acc)
    txt(s, f"{i+1}. {lbl}", 0.65, cy9+0.06, 1.9, 0.55, size=13, bold=True, color=acc)
    txt(s, desc, 2.7, cy9+0.18, 9.9, 0.45, size=13, color=TEXT_HI)

txt(s, "Output: connectome.json — drop directly into Brain-Viz via drag-and-drop",
    0.65, 7.1, 12.0, 0.4, size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 10  —  VISUALIZATION FRAMEWORK
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Visualization Framework: Three.js + D3.js", ACCENT)

for cx10, title10, acc10, items10 in [
    (0.45, "3D Spatial View  (Three.js r128 WebGL)", ACCENT, [
        "InstancedMesh for 100–400 brain region nodes (GPU instancing)",
        "Vertex-displaced brain shell with 5 anatomical lobe colours",
        "LineSegments for up to 700 structural connectivity edges",
        "Separate dim-mesh renders filtered-out networks at low opacity",
        "OrbitControls: pan, zoom, rotate with damping",
        "Raycaster for hover & click hit-testing on instanced geometry",
        "Animated path particles: CanvasTexture glow dots via PointsMaterial",
        "Coronal / Axial / Sagittal preset camera views",
    ]),
    (6.95, "2D Topological View  (D3.js v7)", TEAL, [
        "Circular layout: regions arranged by community on a ring",
        "Bezier-bundled edges (β=0.82) reduce visual clutter",
        "Community arcs + labels around the outer ring",
        "Chord diagram mode: shows community-level flow matrix",
        "Toggle button switches between circular ↔ chord instantly",
        "Hover highlights all edges of the hovered node",
        "Double-click node triggers pathfinding across both views",
        "Brushing & linking — store events sync both views",
    ]),
]:
    card(s, cx10, 1.2, 6.1, 5.9, BG_CARD)
    rect(s, cx10, 1.2, 6.1, 0.07, acc10)
    txt(s, title10, cx10+0.2, 1.28, 5.7, 0.5, size=15, bold=True, color=acc10)
    bullet_block(s, items10, cx10+0.2, 1.88, 5.7, 5.0, size=12.5, accent=acc10, spacing=0.65)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 11  —  KEY FEATURES OVERVIEW
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Key Features at a Glance", GREEN)

features = [
    ("🧠",  "3D Brain Shell",        ACCENT, "Vertex-coloured lobe anatomy\nrealistic cortical shape"),
    ("🔗",  "Shortest Pathway",       GOLD,   "Dijkstra pathfinding with\nanimated particle flow"),
    ("⚖",   "Dataset Comparison",    GREEN,  "Side-by-side healthy vs.\ndisease connectome diff"),
    ("💡",  "Region Info Panel",     TEAL,   "Clinical description for\nevery brain region clicked"),
    ("🎯",  "Community Filter",      PINK,   "Isolate networks with\nopacity-dimmed remainder"),
    ("📊",  "Network Metrics",       ORANGE, "Global efficiency, density,\nbetweenness centrality"),
    ("🔍",  "Region Search",         ACCENT, "Type-to-search with\nhover & pathfinding"),
    ("📁",  "Multi-dataset",         TEAL,   "Load unlimited files;\nhot-swap datasets"),
    ("🌙",  "Dark / Light mode",     GOLD,   "Toggle with WebGL\nbackground sync"),
    ("⬇",   "PNG Export",            GREEN,  "Download current 3D view\nas timestamped PNG"),
    ("🎼",  "Chord Diagram",         PINK,   "Community-level flow\nD3.js chord toggle"),
    ("💊",  "ABIDE Integration",     ORANGE, "Real autism vs. control\nfMRI connectivity data"),
]
for i,(icon,name,acc,desc) in enumerate(features):
    row,col = divmod(i,4)
    cx11 = 0.45 + col*3.22
    cy11 = 1.2 + row*2.05
    card(s, cx11, cy11, 3.0, 1.82, BG_CARD)
    rect(s, cx11, cy11, 3.0, 0.07, acc)
    txt(s, icon, cx11+0.15, cy11+0.12, 0.55, 0.55, size=22)
    txt(s, name, cx11+0.75, cy11+0.12, 2.1, 0.45, size=13, bold=True, color=acc)
    txt(s, desc, cx11+0.15, cy11+0.68, 2.75, 0.95, size=11.5, color=TEXT_HI)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 12  —  3D SPATIAL VIEW RESULTS
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Results: 3D Spatial View — Schaefer-100 Atlas", ACCENT)

card(s, 0.45, 1.2, 8.5, 5.9, BG_MID)
txt(s, "3D Brain Visualization", 0.65, 1.3, 8.1, 0.4, size=16, bold=True, color=ACCENT)
# Simulate brain visualization with shapes
sh_brain = s.shapes.add_shape(9, Inches(2.0), Inches(1.85), Inches(5.5), Inches(4.5))
sh_brain.fill.solid(); sh_brain.fill.fore_color.rgb = RGBColor(0x1C,0x28,0x44)
sh_brain.line.color.rgb = RGBColor(0x3A,0x86,0xFF)

txt(s, "🧠  Schaefer-100 parcellation\n"
       "with anatomical lobe colouring:\n\n"
       "  ■  Frontal (Blue)\n"
       "  ■  Parietal (Teal)\n"
       "  ■  Temporal (Orange)\n"
       "  ■  Occipital (Purple)\n"
       "  ■  Insula (Gold)\n\n"
       "OrbitControls for free rotation\n"
       "Front / Top / Side presets",
    2.2, 2.1, 5.0, 4.0, size=14, color=TEXT_HI)

# Right annotations
card(s, 9.2, 1.2, 4.05, 1.62, BG_CARD)
rect(s, 9.2, 1.2, 4.05, 0.07, ACCENT)
txt(s, "Node encoding", 9.4, 1.28, 3.6, 0.4, size=13, bold=True, color=ACCENT)
for lbl, desc in [("Size", "Degree centrality"), ("Colour", "Functional network"), ("Glow", "Path / selected")]:
    pass
txt(s, "Size → degree centrality\nColour → functional network\nGlow → path / selected node",
    9.4, 1.75, 3.6, 1.0, size=12, color=TEXT_HI)

card(s, 9.2, 3.05, 4.05, 1.62, BG_CARD)
rect(s, 9.2, 3.05, 4.05, 0.07, GOLD)
txt(s, "Edge encoding", 9.4, 3.13, 3.6, 0.4, size=13, bold=True, color=GOLD)
txt(s, "Opacity → weight threshold\nColour → blended community\nGold → shortest path\nDim → filtered-out networks",
    9.4, 3.60, 3.6, 1.0, size=12, color=TEXT_HI)

card(s, 9.2, 4.9, 4.05, 2.2, BG_CARD)
rect(s, 9.2, 4.9, 4.05, 0.07, TEAL)
txt(s, "Performance", 9.4, 4.98, 3.6, 0.4, size=13, bold=True, color=TEAL)
txt(s, "100 nodes · 700 edges\nSteady 60 FPS on mid-range GPU\nInstancedMesh: 1 draw call for all nodes\nLineSegments: 1 draw call for all edges\nParticle system: 14 dots at 60 FPS",
    9.4, 5.45, 3.6, 1.55, size=11.5, color=TEXT_HI)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 13  —  2D TOPOLOGY & CHORD
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Results: 2D Topological View — Circular & Chord Diagrams", TEAL)

for cx13, title13, acc13, items13 in [
    (0.45, "Circular Layout", TEAL, [
        "100 nodes arranged by community on a ring",
        "Bezier-bundled edges reduce visual clutter",
        "Community arcs + network labels around perimeter",
        "Hover: highlights all edges of the hovered node",
        "Edge opacity proportional to connection weight",
    ]),
    (6.95, "Chord Diagram Mode", PINK, [
        "Aggregates edges into community × community flow matrix",
        "Arc width proportional to community size",
        "Ribbon width proportional to inter-community connectivity",
        "Hover arc: dims all ribbons except that community's",
        "Toggle button switches mode without page reload",
    ]),
]:
    card(s, cx13, 1.2, 6.1, 3.0, BG_CARD)
    rect(s, cx13, 1.2, 6.1, 0.07, acc13)
    txt(s, title13, cx13+0.2, 1.28, 5.7, 0.45, size=17, bold=True, color=acc13)
    bullet_block(s, items13, cx13+0.2, 1.82, 5.7, 5.0, size=13, accent=acc13, spacing=0.52)

card(s, 0.45, 4.4, 12.4, 2.7, BG_MID)
txt(s, "Brushing & Linking", 0.65, 4.5, 4.0, 0.4, size=15, bold=True, color=GREEN)
link_items = [
    "Hover node in 3D view → same node highlighted in 2D view",
    "Click node in 2D view → pathfinding source set in both views",
    "Community filter toggle → both 3D and 2D views update simultaneously via BrainStore event bus",
    "Weight threshold slider → both views update visible edges in real-time",
]
bullet_block(s, link_items, 0.65, 5.0, 12.0, 2.0, size=13, accent=GREEN, spacing=0.52)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 14  —  DATASET COMPARISON FEATURE
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Results: Dataset Comparison — Structural Diff View", GREEN)

card(s, 0.45, 1.2, 12.4, 2.6, BG_MID)
txt(s, "How Comparison Works", 0.65, 1.28, 5.0, 0.4, size=16, bold=True, color=GREEN)
cmp_items = [
    "Load two connectome.json files on the landing page",
    "Edge diff: classified as 'Only in A' / 'Only in B' / 'Shared' by normalised edge key",
    "Two independent Three.js renderers side-by-side — each with its own scene, camera, OrbitControls",
    "Brushing: hover any node in renderer A → matching label highlighted in renderer B",
]
bullet_block(s, cmp_items, 0.65, 1.75, 12.0, 1.9, size=13.5, accent=GREEN, spacing=0.50)

for cx14, lbl14, col14, items14 in [
    (0.45, "Only in A", RGBColor(0x00,0xE6,0x76), [
        "Edges present in dataset A but absent in B",
        "Rendered in bright GREEN",
        "Indicates stronger or unique connectivity in A",
    ]),
    (4.6,  "Shared",    RGBColor(0xBB,0xBB,0xBB), [
        "Edges present in both datasets",
        "Rendered in WHITE / grey",
        "Core structural backbone preserved across groups",
    ]),
    (8.75, "Only in B",  RGBColor(0xFF,0x52,0x52), [
        "Edges present in dataset B but absent in A",
        "Rendered in bright RED",
        "Indicates stronger or unique connectivity in B",
    ]),
]:
    card(s, cx14, 4.0, 4.0, 2.15, BG_CARD)
    rect(s, cx14, 4.0, 4.0, 0.07, col14)
    txt(s, lbl14, cx14+0.2, 4.1, 3.6, 0.45, size=15, bold=True, color=col14)
    bullet_block(s, items14, cx14+0.2, 4.62, 3.6, 1.5, size=12.5, accent=col14, spacing=0.5)

card(s, 0.45, 6.3, 12.4, 0.9, BG_CARD)
txt(s, "Bottom stats bar: ABIDE autism vs. control — 630 (73.8%) unique to A  ·  70 (8.2%) shared  ·  154 (18.0%) unique to B",
    0.65, 6.38, 12.0, 0.65, size=13.5, color=TEXT_HI, align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 15  —  REGION INFO PANEL
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Results: Region Info Panel — Clinical Knowledge Integration", GOLD)

card(s, 0.45, 1.2, 7.3, 5.9, BG_CARD)
txt(s, "What the Panel Shows", 0.65, 1.3, 6.9, 0.45, size=18, bold=True, color=GOLD)
panel_items = [
    "Lobe badge (Frontal/Parietal/Temporal/Occipital/Insula/Central) with colour",
    "Functional network badge (Default Mode / Visual / Somatomotor …)",
    "Plain-English description of what the region does",
    "Clinical relevance — what damage here causes in patients",
    "Network role — which large-scale network this region belongs to",
    "Connectivity metrics: degree, % of regions connected, betweenness centrality",
    "Hub classification: Major Hub / Secondary Hub / Local Node",
    "Layman analogy: 'This region acts like a major airport…'",
]
bullet_block(s, panel_items, 0.65, 1.85, 7.0, 5.0, size=13.5, accent=GOLD, spacing=0.65)

card(s, 7.95, 1.2, 4.9, 5.9, BG_MID)
txt(s, "Example: L-Precuneus", 8.15, 1.3, 4.5, 0.45, size=15, bold=True, color=GOLD)
example_lines = [
    ("Lobe:",       "Parietal", TEAL),
    ("Network:",    "Default Mode", ACCENT),
    ("Function:",   "Self-referential thought,\nvisuospatial imagery,\nepisodic memory", TEXT_HI),
    ("Clinical:",   "Atrophy linked to\nAlzheimer's & depression", PINK),
    ("Degree:",     "12 connections  (12.1%)", GOLD),
    ("Hub Class:",  "Secondary Hub", GREEN),
]
cy15 = 1.85
for lbl,val,col in example_lines:
    txt(s, lbl, 8.15, cy15, 1.5, 0.38, size=11.5, bold=True, color=TEXT_MID)
    txt(s, val, 9.75, cy15, 3.0, 0.70, size=12.5, color=col)
    cy15 += 0.72

card(s, 7.95, 6.3, 4.9, 0.7, RGBColor(0x1C,0x28,0x44))
txt(s, "💡  In plain English: This region acts like a regional hub — fairly well-connected, important for local coordination in self-awareness and memory.",
    8.1, 6.35, 4.6, 0.58, size=11, color=TEXT_HI, italic=True)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 16  —  REAL DATA: ABIDE RESULTS
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Results: Real Disease Data — ABIDE Autism vs. Control", PINK)

txt(s, "Dataset: ABIDE (Autism Brain Imaging Data Exchange)  ·  CPAC pipeline  ·  5 subjects per group  ·  Schaefer-100 atlas",
    0.45, 1.15, 12.4, 0.38, size=13, color=TEXT_MID, italic=True)

for cx16, grp, acc16, stats in [
    (0.45,  "Control (Neurotypical)", GREEN, [
        ("700 edges",        "after threshold 0.05", GOLD),
        ("Density 0.1414",   "denser connectivity", TEAL),
        ("14.0 avg degree",  "highly connected hubs", ACCENT),
        ("Strong DMN",       "Default Mode backbone intact", GREEN),
    ]),
    (6.9, "ASD (Autism Spectrum)", PINK, [
        ("224 edges",        "after threshold 0.05", GOLD),
        ("Density 0.0453",   "sparser connectivity", TEAL),
        ("4.5 avg degree",   "reduced hub connectivity", ACCENT),
        ("Weak DMN",         "Default Mode disrupted", PINK),
    ]),
]:
    card(s, cx16, 1.6, 5.95, 2.3, BG_CARD)
    rect(s, cx16, 1.6, 5.95, 0.07, acc16)
    txt(s, grp, cx16+0.2, 1.68, 5.5, 0.45, size=16, bold=True, color=acc16)
    for j,(val,desc,col) in enumerate(stats):
        stat_box(s, val, desc, cx16+0.2 + j*1.38, 2.18, 1.24, col)

card(s, 0.45, 4.1, 12.4, 3.0, BG_MID)
txt(s, "Key Findings from Visual Comparison", 0.65, 4.18, 8.0, 0.45, size=16, bold=True, color=PINK)
findings = [
    "Control brains show significantly denser connectivity (700 vs 224 edges; 3× more structural pathways)",
    "ASD group exhibits reduced long-range connectivity especially between frontal and parietal regions",
    "Default Mode Network (DMN) shows 68% fewer edges in autism group — consistent with published literature",
    "Dorsal Attention Network connections are preserved in ASD (no significant difference in edge count)",
    "Hub analysis: top 3 hub regions shift in autism — thalamus becomes more prominent, precuneus less central",
    "Structural diff overlay: 630 edges unique to controls vs 154 unique to ASD (73.8% vs 18.0%)",
]
bullet_block(s, findings, 0.65, 4.7, 12.0, 2.5, size=13, accent=PINK, spacing=0.43)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 17  —  TECHNICAL PERFORMANCE
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Technical Performance & Implementation Metrics", TEAL)

metrics = [
    ("60 FPS",      "3D rendering\n100 nodes, 700 edges",          ACCENT),
    ("<5 ms",       "Dijkstra runtime\n100-node graph",             GOLD),
    ("~2 MB",       "Atlas download\nSchaefer-100 (cached)",        TEAL),
    ("0 install",   "Browser-native\nno plugins needed",            GREEN),
    ("912 lines",   "main.js\nall features wired",                  PINK),
    ("857 lines",   "brain3d.js\nThree.js WebGL engine",           ORANGE),
]
for i,(val,lbl,acc) in enumerate(metrics):
    row17, col17 = divmod(i, 3)
    stat_box(s, val, lbl, 0.45 + col17*4.18, 1.25 + row17*2.0, 3.95, acc)

card(s, 0.45, 5.5, 12.4, 1.65, BG_CARD)
txt(s, "Technology Stack", 0.65, 5.58, 4.0, 0.4, size=14, bold=True, color=TEXT_HI)
tech_items = [
    "Three.js r128  ·  WebGL 3D rendering with InstancedMesh",
    "D3.js v7  ·  SVG-based 2D chord + circular topology views",
    "nilearn (Python)  ·  Atlas download + fMRI connectivity computation",
    "python-pptx / pptxgenjs  ·  Presentation generation",
]
bullet_block(s, tech_items, 0.65, 6.0, 12.0, 1.3, size=12.5, spacing=0.38)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 18  —  FINDINGS & INSIGHTS
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Findings & Visual Insights", GREEN)

categories = [
    ("Technical Findings",  ACCENT, [
        "InstancedMesh + LineSegments: 2 draw calls render the full brain — enables 60 FPS on integrated GPU",
        "Dim-mesh pattern: two separate InstancedMesh objects (active + dim) give opacity-per-group without shader complexity",
        "Circular particle texture (CanvasTexture radial gradient) eliminates the PointsMaterial square artefact",
        "MNI axis remapping (x→X, z→Y, y→Z) distributes nodes in anatomically correct brain shape",
    ]),
    ("Neuroscience Insights", GREEN, [
        "ABIDE fMRI data reveals 3× denser connectivity in neurotypical vs. ASD cohort (700 vs 224 edges)",
        "Default Mode Network (precuneus, mPFC, PCC) exhibits greatest connectivity reduction in autism",
        "Hub regions shift: thalamus centrality ↑ in ASD; precuneus centrality ↓ — matches published findings",
        "Global efficiency in fMRI: 0.041 (autism) vs 0.287 (healthy nilearn) — significant integration drop",
    ]),
    ("UX Insights", GOLD, [
        "Brushing & linking across views is the single most impactful interaction for exploratory analysis",
        "Region info panel in sidebar (not floating) prevents occluding the 2D view during investigation",
        "Landing page mode selection (visualize vs. compare) reduces cognitive load for first-time users",
        "Layperson plain-English descriptions ('major airport analogy') validated by non-expert test users",
    ]),
]
cy18 = 1.2
for title18, acc18, items18 in categories:
    card(s, 0.45, cy18, 12.4, 1.82, BG_CARD)
    rect(s, 0.45, cy18, 0.07, 1.82, acc18)
    txt(s, title18, 0.65, cy18+0.08, 3.5, 0.4, size=14, bold=True, color=acc18)
    for j,item in enumerate(items18):
        txt(s, f"·  {item}", 4.3, cy18+0.1+j*0.43, 8.3, 0.4, size=12, color=TEXT_HI)
    cy18 += 2.0

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 19  —  CONCLUSION & FUTURE WORK
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)
title_bar(s, "Conclusion & Future Work", ACCENT)

card(s, 0.45, 1.2, 6.0, 5.9, BG_CARD)
txt(s, "What We Achieved ✅", 0.65, 1.3, 5.6, 0.45, size=18, bold=True, color=GREEN)
achieved = [
    "Fully functional web-based connectome dashboard",
    "3D WebGL brain with anatomical lobe colouring",
    "2D D3.js dual-layout (circular + chord diagram)",
    "Dijkstra pathfinding with animated path particles",
    "Louvain community detection + network filtering",
    "Real ABIDE dataset integration (autism vs. control)",
    "Dataset comparison with structural diff overlay",
    "Region info panel with clinical knowledge base",
    "Brushing & linking across all views",
    "Zero-install, browser-native deployment",
]
bullet_block(s, achieved, 0.65, 1.85, 5.7, 5.0, size=13, accent=GREEN, spacing=0.50)

card(s, 6.65, 1.2, 6.1, 5.9, BG_CARD)
txt(s, "Future Work 🔭", 6.85, 1.3, 5.7, 0.45, size=18, bold=True, color=PINK)
future = [
    "Actual 3D brain mesh overlay (FreeSurfer .obj) as glass shell",
    "Group-level statistics (t-test, effect size) in comparison view",
    "Multi-subject averaging (10+ subjects per condition)",
    "Time-series fMRI animation — watch connectivity evolve",
    "ADNI Alzheimer's dataset integration via DTI pipeline",
    "Machine learning classifier: predict diagnosis from connectivity",
    "Export connectivity matrix as CSV / NetworkX GraphML",
    "WebXR mode: explore the brain in VR headset",
    "Accessibility: screen-reader support, colour-blind palette",
    "Mobile-responsive layout for tablet & phone viewing",
]
bullet_block(s, future, 6.85, 1.85, 5.7, 5.0, size=13, accent=PINK, spacing=0.50)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 20  —  THANK YOU
# ════════════════════════════════════════════════════════════════════════════════
s = blank_slide(); bg(s, BG_DARK)

# Full bleed left panel
rect(s, 0, 0, 6.5, 7.5, BG_MID)
rect(s, 0, 0, 0.12, 7.5, ACCENT)

# Decorative dots
import random; random.seed(99)
for _ in range(30):
    dx = random.uniform(0.3, 5.9)
    dy = random.uniform(0.3, 7.0)
    dr = random.uniform(0.05, 0.20)
    dc = random.choice([ACCENT,TEAL,GOLD,PINK,GREEN])
    sh = s.shapes.add_shape(9, Inches(dx-dr), Inches(dy-dr), Inches(dr*2), Inches(dr*2))
    sh.fill.solid(); sh.fill.fore_color.rgb = dc
    sh.line.fill.background()

txt(s, "Thank You", 0.5, 2.0, 5.6, 1.5, size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Questions & Discussion", 0.5, 3.5, 5.6, 0.7, size=22, color=TEAL, italic=True, align=PP_ALIGN.CENTER)

rect(s, 0.7, 4.35, 5.0, 0.04, ACCENT)

txt(s, "Nishant · 2025EET2734", 0.5, 4.55, 5.6, 0.4, size=15, color=TEXT_MID, align=PP_ALIGN.CENTER)
txt(s, "Aryan · 2025EET2486",   0.5, 4.95, 5.6, 0.4, size=15, color=TEXT_MID, align=PP_ALIGN.CENTER)
txt(s, "EET — Data Visualization  |  2025", 0.5, 5.55, 5.6, 0.35, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Right panel — key links / references
card(s, 6.8, 0.9, 6.1, 5.8, BG_CARD)
txt(s, "Project Summary", 7.0, 1.0, 5.7, 0.45, size=18, bold=True, color=ACCENT)
summary = [
    ("Tool Stack",    "Three.js r128  ·  D3.js v7  ·  Python nilearn"),
    ("Algorithms",    "Dijkstra  ·  Louvain  ·  Brandes betweenness"),
    ("Data",          "HCP  ·  ABIDE  ·  Schaefer-100 atlas"),
    ("Key Feature",   "3D brain  +  2D chord  +  disease comparison"),
    ("Lines of Code", "~3500 lines across 6 JavaScript files"),
    ("Dataset Size",  "100 nodes · up to 700 edges per connectome"),
]
cy20 = 1.6
for lbl,val in summary:
    txt(s, lbl+":", 7.0, cy20, 2.1, 0.38, size=12.5, bold=True, color=TEXT_MID)
    txt(s, val,     9.2, cy20, 3.5, 0.38, size=12.5, color=TEXT_HI)
    cy20 += 0.72

card(s, 6.8, 6.9, 6.1, 0.6, RGBColor(0x1C,0x28,0x44))
txt(s, "localhost:8080  ·  python -m http.server 8080",
    7.0, 6.96, 5.7, 0.4, size=12, color=TEAL, italic=True, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/nisha/Downloads/brain-viz/BrainViz_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")